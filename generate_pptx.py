from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Layouts: 0 is Title, 1 is Title and Content
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bridging the Silence: Real-Time AI Sign Language Translation"
    subtitle.text = "Empowering communication through Machine Learning and Computer Vision\nPresented by Ethiopia Chan & Betelehem Addisu"

    # 2. The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Communication Barrier"
    content = slide.placeholders[1]
    content.text = ("• Over 430 million people worldwide experience disabling hearing loss.\n"
                    "• Sign language is a rich, complex language, but very few hearing people understand it.\n"
                    "• Traditional translators are expensive and not available 24/7.\n"
                    "• Result: Social isolation and unequal access to daily services.")

    # 3. Our Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "An Accessible AI Interpreter"
    content = slide.placeholders[1]
    content.text = ("• Real-Time Translation: Translates signs into text instantly.\n"
                    "• Voice Output: Converts text into spoken audio automatically.\n"
                    "• Cost-Effective: Runs on standard webcams—no expensive sensory gloves.\n"
                    "• User-Friendly: Interactive dashboard designed for everyone.")

    # 4. How It Works (Technical)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Under the Hood (Technology)"
    content = slide.placeholders[1]
    content.text = ("• Data Capture: OpenCV captures live video frames.\n"
                    "• Landmark Detection: MediaPipe maps 21 3D coordinates on the hand.\n"
                    "• Machine Learning: Custom Convolutional Neural Network (CNN) for classification.\n"
                    "• Natural Language Processing: Smart spelling suggestions and AI voice.")

    # 5. STEAM Integration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "A True STEAM Project"
    content = slide.placeholders[1]
    content.text = ("• Science: Human anatomy and linguistic gestures.\n"
                    "• Technology: Computer Vision, AI, and TTS engines.\n"
                    "• Engineering: Software architecture and UI/UX design.\n"
                    "• Arts: Designing an intuitive and accessible dashboard.\n"
                    "• Mathematics: Geometry for hand joint distances and CNN Linear Algebra.")

    # 6. Meet the Team
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Minds Behind the Code"
    content = slide.placeholders[1]
    content.text = ("Ethiopia Chan (Founder)\n"
                    "• Software Engineer & Junior ML Student\n"
                    "• 12th Grade Student at GSS\n\n"
                    "Betelehem Addisu (Co-Founder)\n"
                    "• Data Management & Model Analyzer\n"
                    "• 12th Grade Student at ELAY")

    # 7. Future Vision
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Looking Forward"
    content = slide.placeholders[1]
    content.text = ("• Mobile App: Porting the model to iOS and Android.\n"
                    "• Dynamic Words: Expanding to moving signs and gestures.\n"
                    "• Global Languages: Adapting for Ethiopian Sign Language (EthSL).\n"
                    "• Goal: Eradicating the communication barrier entirely.")

    # Save
    file_path = "Sign_Language_Interpreter_Presentation.pptx"
    prs.save(file_path)
    print(f"Presentation saved to {file_path}")

if __name__ == "__main__":
    create_presentation()
